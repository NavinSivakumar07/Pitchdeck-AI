"""
Content Agent Module
Implements CrewAI agent that generates PowerPoint presentations based on knowledge analysis and research data.
"""

import os
import logging
import json
from typing import Dict, Any
from datetime import datetime
from crewai import Agent, Task, Crew, Process, LLM
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from dotenv import load_dotenv

# Load environment variables
load_dotenv()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class ContentAgent:
    """Content Agent that generates PowerPoint presentations based on analysis."""
    
    def __init__(self):
        """Initialize the Content Agent with LLM."""
        # Configure LLM
        self.llm = LLM(
            model="gpt-4o-mini",
            api_key=os.getenv("OPENAI_API_KEY"),
            temperature=0.4,
            max_tokens=3000
        )
        
        # Initialize the agent
        self.agent = Agent(
            role="Pitch Deck Content Creator & Presentation Designer",
            goal="Create compelling pitch deck content and PowerPoint presentations based on research and knowledge analysis",
            backstory="""You are an expert pitch deck designer and content creator with extensive experience 
            in creating successful startup presentations. You understand what investors look for and how to 
            structure compelling narratives that secure funding.""",
            tools=[],
            allow_delegation=False,
            verbose=True,
            llm=self.llm
        )
    
    def generate_pitch_content(self, startup_data: Dict[str, Any], research_output: str, knowledge_analysis: str) -> str:
        """
        Generate structured pitch deck content based on all available data.
        
        Args:
            startup_data: Original startup information
            research_output: Research agent output
            knowledge_analysis: Knowledge agent analysis
            
        Returns:
            Structured pitch deck content
        """
        try:
            company_name = startup_data.get('startup_name', 'Unknown Company')
            
            # Create content generation task
            content_task = Task(
                description=f"""
                Create comprehensive pitch deck content for '{company_name}' based on the provided information.
                
                Startup Data:
                {json.dumps(startup_data, indent=2)}
                
                Research Output:
                {research_output}
                
                Knowledge Analysis:
                {knowledge_analysis}
                
                Generate content for a 12-slide pitch deck with the following structure:
                1. Title Slide - Company name, tagline, founder info
                2. Problem Statement - Clear problem definition with market evidence
                3. Solution - Your unique solution and value proposition
                4. Market Opportunity - Market size, growth, and target segments
                5. Product Demo - Key features and benefits
                6. Business Model - Revenue streams and pricing strategy
                7. Traction - Current progress, metrics, and achievements
                8. Competition - Competitive landscape and differentiation
                9. Marketing Strategy - Customer acquisition and growth plan
                10. Team - Founder and key team member profiles
                11. Financial Projections - Revenue forecasts and funding use
                12. Funding Ask - Investment amount and terms
                
                For each slide, provide:
                - Slide title
                - Key bullet points (3-5 per slide)
                - Supporting data/statistics where relevant
                - Visual suggestions (charts, images, diagrams)
                - Speaker notes with additional context
                
                Base your content on successful patterns from the knowledge analysis and incorporate 
                insights from the research output.
                """,
                expected_output="""
                A clean, readable pitch deck format with each slide clearly structured as:

                Slide X: [Title]
                [Content with bullet points and key information]

                Each slide should be investor-ready with compelling content that tells a cohesive story.
                Format should be clean and easy to read, not JSON.
                """,
                agent=self.agent
            )
            
            # Create and run crew
            content_crew = Crew(
                agents=[self.agent],
                tasks=[content_task],
                process=Process.sequential,
                verbose=True,
                full_output=True
            )
            
            logger.info("Generating pitch deck content...")
            result = content_crew.kickoff()

            # Generate clean formatted output with research and knowledge insights
            clean_output = self.format_clean_pitch_output(startup_data, research_output, knowledge_analysis)

            # Save content result
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            content_file = f"pitch_content_{company_name}_{timestamp}.txt"

            with open(content_file, "w", encoding="utf-8") as f:
                f.write(f"# Pitch Deck Content: {company_name}\n")
                f.write("=" * 60 + "\n\n")
                f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n\n")
                f.write(clean_output)

            logger.info(f"Pitch content generated and saved to {content_file}")
            return clean_output
            
        except Exception as e:
            logger.error(f"Error generating pitch content: {e}")
            import traceback
            logger.error(traceback.format_exc())
            return f"Error generating content: {str(e)}"
    
    def create_powerpoint_presentation(self, startup_data: Dict[str, Any], pitch_content: str) -> str:
        """
        Create a PowerPoint presentation from the generated content.
        
        Args:
            startup_data: Original startup information
            pitch_content: Generated pitch deck content
            
        Returns:
            Path to the created PowerPoint file
        """
        try:
            company_name = startup_data.get('startup_name', 'Unknown Company')
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            ppt_filename = f"pitch_deck_{company_name}_{timestamp}.pptx"
            
            # Create presentation
            prs = Presentation()
            
            # Define slide layouts and styling
            title_slide_layout = prs.slide_layouts[0]  # Title slide
            content_slide_layout = prs.slide_layouts[1]  # Title and content
            
            # Parse content (simplified parsing - in production, you'd want more robust parsing)
            slides_data = self._parse_pitch_content(pitch_content, startup_data)
            
            # Create slides
            for slide_data in slides_data:
                if slide_data['slide_number'] == 1:
                    # Title slide
                    slide = prs.slides.add_slide(title_slide_layout)
                    title = slide.shapes.title
                    subtitle = slide.placeholders[1]
                    
                    title.text = slide_data['slide_title']
                    subtitle.text = f"Founder: {startup_data.get('founder_name', 'N/A')}\n{datetime.now().strftime('%B %Y')}"
                else:
                    # Content slide
                    slide = prs.slides.add_slide(content_slide_layout)
                    title = slide.shapes.title
                    content = slide.placeholders[1]
                    
                    title.text = slide_data['slide_title']
                    
                    # Add bullet points
                    text_frame = content.text_frame
                    text_frame.clear()
                    
                    for i, bullet in enumerate(slide_data['bullet_points']):
                        if i == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = bullet
                        p.level = 0
                        p.font.size = Pt(18)
            
            # Save presentation
            prs.save(ppt_filename)
            logger.info(f"PowerPoint presentation created: {ppt_filename}")
            
            return ppt_filename
            
        except Exception as e:
            logger.error(f"Error creating PowerPoint presentation: {e}")
            import traceback
            logger.error(traceback.format_exc())
            return f"Error creating presentation: {str(e)}"
    
    def _parse_pitch_content(self, pitch_content: str, startup_data: Dict[str, Any]) -> list:
        """
        Parse the generated pitch content into structured slide data.
        
        Args:
            pitch_content: Raw pitch content from agent
            startup_data: Original startup data
            
        Returns:
            List of slide data dictionaries
        """
        # This is a simplified parser - in production, you'd want more robust parsing
        slides = []
        
        # Default slide structure based on startup data
        default_slides = [
            {
                'slide_number': 1,
                'slide_title': startup_data.get('startup_name', 'Company Name'),
                'bullet_points': [
                    startup_data.get('vision_statement', 'Vision statement'),
                    f"Founder: {startup_data.get('founder_name', 'N/A')}",
                    f"Industry: {startup_data.get('industry_type', 'N/A')}"
                ]
            },
            {
                'slide_number': 2,
                'slide_title': 'Problem',
                'bullet_points': [
                    startup_data.get('key_problem_solved', 'Problem statement'),
                    'Market pain points',
                    'Current solutions are inadequate'
                ]
            },
            {
                'slide_number': 3,
                'slide_title': 'Solution',
                'bullet_points': [
                    startup_data.get('solution_summary', 'Solution overview'),
                    'Unique value proposition',
                    'Key benefits'
                ]
            },
            {
                'slide_number': 4,
                'slide_title': 'Market Opportunity',
                'bullet_points': [
                    f"Market Size: {startup_data.get('market_size', 'TBD')}",
                    'Target market segments',
                    'Growth potential'
                ]
            },
            {
                'slide_number': 5,
                'slide_title': 'Product',
                'bullet_points': [
                    f"Product: {startup_data.get('product_name', 'Product name')}",
                    'Key features',
                    'User benefits'
                ]
            },
            {
                'slide_number': 6,
                'slide_title': 'Business Model',
                'bullet_points': [
                    startup_data.get('business_model', 'Business model'),
                    startup_data.get('monetization_plan', 'Monetization strategy'),
                    'Revenue streams'
                ]
            },
            {
                'slide_number': 7,
                'slide_title': 'Traction',
                'bullet_points': [
                    startup_data.get('transactions', 'Current traction'),
                    'Key metrics',
                    'Growth milestones'
                ]
            },
            {
                'slide_number': 8,
                'slide_title': 'Competition',
                'bullet_points': [
                    f"Competitors: {startup_data.get('competitors', 'TBD')}",
                    startup_data.get('why_you_win', 'Competitive advantage'),
                    'Market positioning'
                ]
            },
            {
                'slide_number': 9,
                'slide_title': 'Marketing Strategy',
                'bullet_points': [
                    startup_data.get('acquisition_strategy', 'Customer acquisition'),
                    f"Target customers: {startup_data.get('target_customer_profile', 'TBD')}",
                    'Growth channels'
                ]
            },
            {
                'slide_number': 10,
                'slide_title': 'Team',
                'bullet_points': [
                    f"Founder: {startup_data.get('founder_name', 'N/A')}",
                    startup_data.get('founder_bio', 'Founder background'),
                    startup_data.get('team_summary', 'Team overview')
                ]
            },
            {
                'slide_number': 11,
                'slide_title': 'Financial Projections',
                'bullet_points': [
                    'Revenue projections',
                    'Key financial metrics',
                    'Path to profitability'
                ]
            },
            {
                'slide_number': 12,
                'slide_title': 'Funding Ask',
                'bullet_points': [
                    f"Raising: {startup_data.get('funding_amount', 'TBD')}",
                    f"Use of funds: {startup_data.get('use_of_funds_split_percentages', 'TBD')}",
                    'Expected outcomes'
                ]
            }
        ]
        
        return default_slides

    def format_clean_pitch_output(self, startup_data: Dict[str, Any], research_output: str = "", knowledge_analysis: str = "") -> str:
        """
        Generate clean, readable pitch deck format with enhanced content and sourcing.

        Args:
            startup_data: Original startup information
            research_output: Research agent findings
            knowledge_analysis: Knowledge agent analysis

        Returns:
            Clean formatted pitch deck content with sources
        """
        company_name = startup_data.get('startup_name', 'Company Name')
        founder_name = startup_data.get('founder_name', 'Founder Name')
        industry = startup_data.get('industry_type', 'Industry')
        product_name = startup_data.get('product_name', 'Product')
        vision = startup_data.get('vision_statement', 'Vision statement')
        problem = startup_data.get('key_problem_solved', 'Problem statement')
        solution = startup_data.get('solution_summary', 'Solution overview')
        market_size = startup_data.get('market_size', 'Market size')
        business_model = startup_data.get('business_model', 'Business model')
        competitors = startup_data.get('competitors', 'Competitors')
        why_win = startup_data.get('why_you_win', 'Competitive advantage')
        traction = startup_data.get('transactions', 'Current traction')
        team = startup_data.get('team_summary', 'Team overview')
        funding_amount = startup_data.get('funding_amount', 'Funding amount')
        use_of_funds = startup_data.get('use_of_funds_split_percentages', 'Use of funds')
        acquisition_strategy = startup_data.get('acquisition_strategy', 'Customer acquisition strategy')
        target_customers = startup_data.get('target_customer_profile', 'Target customers')
        founder_bio = startup_data.get('founder_bio', 'Founder background')

        # Extract key insights from research and knowledge analysis
        market_insights = self._extract_market_insights(research_output)
        competitive_insights = self._extract_competitive_insights(knowledge_analysis)

        # Format funding amount for display
        try:
            funding_num = float(funding_amount) if funding_amount else 0
            if funding_num >= 1000000:
                funding_display = f"${funding_num/1000000:.1f}M"
            elif funding_num >= 1000:
                funding_display = f"${funding_num/1000:.0f}K"
            else:
                funding_display = f"${funding_num:,.0f}"
        except:
            funding_display = f"${funding_amount}"

        formatted_output = f"""
Slide 1: Title / Company Overview
{company_name}
{vision}
Founder: {founder_name} | Industry: {industry}
📧 contact@{company_name.lower().replace(' ', '')}.com | 🌐 www.{company_name.lower().replace(' ', '')}.com

Slide 2: Problem Statement
🎯 The Challenge:
{problem}

📊 Market Evidence:
• {market_insights.get('problem_validation', 'Significant market pain points identified through research')}
• {market_insights.get('market_size_context', 'Large addressable market with growing demand')}
• Current solutions are inadequate, creating opportunity for innovation

💡 Why Now: Digital transformation and evolving customer needs create urgency for better solutions.

Slide 3: Our Solution
🚀 Introducing {product_name}:
{solution}

✨ Key Differentiators:
• {why_win}
• Proven technology with measurable results
• User-centric design based on customer feedback
• Scalable architecture for future growth

🎯 Value Proposition: We solve the core problem through innovative technology that delivers immediate and long-term value.

Slide 4: Market Opportunity
💰 Total Addressable Market: {market_size}
🎯 Target Segment: {target_customers}

📈 Market Dynamics:
• {market_insights.get('growth_drivers', 'Strong growth driven by digital transformation')}
• {market_insights.get('market_trends', 'Favorable market trends supporting adoption')}
• Expanding customer base with increasing budget allocation

🌍 Geographic Expansion: Initial focus on primary markets with plans for international expansion.

Slide 5: Business Model & Revenue
💼 Revenue Model: {business_model}

💰 Monetization Strategy:
• {startup_data.get('monetization_plan', 'Multiple revenue streams for sustainable growth')}
• Predictable recurring revenue with expansion opportunities
• High customer lifetime value with low churn rates

📊 Unit Economics: Strong margins with scalable cost structure designed for profitability.

Slide 6: Competitive Landscape
🏆 Key Competitors: {competitors}

💪 Our Competitive Advantages:
• {why_win}
• Superior customer experience and support
• Faster time-to-value for customers
• {competitive_insights.get('differentiation', 'Unique positioning in the market')}

🎯 Market Position: {competitive_insights.get('market_position', 'Strong positioning against established players')}

Slide 7: Traction & Validation
📈 Current Traction:
{traction}

🎯 Key Metrics:
• Strong customer satisfaction and retention rates
• Positive unit economics and growth trajectory
• {competitive_insights.get('validation_points', 'Market validation through customer adoption')}

🚀 Momentum: Accelerating growth with increasing market recognition and customer demand.

Slide 8: Leadership Team
👨‍💼 {founder_name} - Founder & CEO
{founder_bio}

👥 Core Team:
{team}

🏆 Advisory Board: Industry experts and successful entrepreneurs providing strategic guidance.

💡 Why This Team Wins: Deep domain expertise, proven execution track record, and complementary skill sets.

Slide 9: Financial Projections & Use of Funds
💰 Funding Ask: {funding_display}

📊 Use of Funds:
{use_of_funds}

📈 Financial Projections:
• 18-24 month runway to achieve key milestones
• Path to profitability with strong unit economics
• Projected 3-5x revenue growth over next 24 months

🎯 Key Milestones: Product development, market expansion, and team scaling.

Slide 10: Go-To-Market Strategy
🎯 Customer Acquisition:
{acquisition_strategy}

📈 Growth Strategy:
• Multi-channel approach with focus on highest ROI channels
• Strategic partnerships for market acceleration
• Content marketing and thought leadership
• Referral programs and customer success initiatives

🌍 Expansion Plan: Geographic and vertical market expansion based on initial success.

Slide 11: Investment Opportunity
💰 We're raising {funding_display} to:
• Accelerate product development and innovation
• Scale go-to-market efforts and customer acquisition
• Expand team with key hires in engineering and sales
• Capture market opportunity and achieve market leadership

🚀 Expected Outcomes:
• 10x revenue growth potential
• Market leadership position in {industry.lower()}
• Strong ROI for investors with clear exit opportunities

Slide 12: Next Steps & Contact
🙏 Thank you for your time and consideration.

🤝 Let's partner to transform the {industry.lower()} industry and create significant value together.

📞 Ready to discuss further:
📧 {founder_name.lower().replace(' ', '.')}@{company_name.lower().replace(' ', '')}.com
📱 [Contact number]
🌐 www.{company_name.lower().replace(' ', '')}.com
💼 LinkedIn: /company/{company_name.lower().replace(' ', '-')}

---
📊 Sources: Market research, competitive analysis, and industry reports
🤖 Generated by AI Pitch Deck Generator with RAG technology
"""

        return formatted_output.strip()

    def _extract_market_insights(self, research_output: str) -> Dict[str, str]:
        """Extract key market insights from research output."""
        insights = {
            'problem_validation': 'Market research confirms significant pain points',
            'market_size_context': 'Large and growing addressable market',
            'growth_drivers': 'Digital transformation driving market expansion',
            'market_trends': 'Favorable trends supporting solution adoption'
        }

        if research_output:
            # Extract specific insights from research
            if 'billion' in research_output.lower():
                insights['market_size_context'] = 'Multi-billion dollar market opportunity'
            if 'growth' in research_output.lower() and '%' in research_output:
                insights['growth_drivers'] = 'Strong market growth with double-digit CAGR'
            if 'trend' in research_output.lower():
                insights['market_trends'] = 'Multiple favorable market trends identified'

        return insights

    def _extract_competitive_insights(self, knowledge_analysis: str) -> Dict[str, str]:
        """Extract competitive insights from knowledge analysis."""
        insights = {
            'differentiation': 'Clear differentiation from existing solutions',
            'market_position': 'Strong competitive positioning',
            'validation_points': 'Validated approach based on successful patterns'
        }

        if knowledge_analysis:
            # Extract specific competitive insights
            if 'competitive advantage' in knowledge_analysis.lower():
                insights['differentiation'] = 'Proven competitive advantages identified'
            if 'market' in knowledge_analysis.lower() and 'position' in knowledge_analysis.lower():
                insights['market_position'] = 'Favorable market positioning vs competitors'
            if 'successful' in knowledge_analysis.lower():
                insights['validation_points'] = 'Success patterns validated through database analysis'

        return insights
